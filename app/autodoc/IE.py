""" 
    Modulo del AutodocIntegrates para generar
    las presentaciones automaticas pptx de un
    proyecto en:
    - Fluid
    - Bancolombia
"""
#third party includes
from pptx import Presentation

class Bancolombia:

    presentation = None
    template_path = "/var/www/fluid-integrates/app/autodoc/templates/bancolombia.pptx"
    result_path = "/var/www/fluid-integrates/app/autodoc/results/:project.pptx"
    critical_finding_slide = 1
    moderate_finding_slide = 10
    tolerable_finding_slide = 27
    MAX_CRITICAL_SLIDE = 10
    MAX_MODERATE_SLIDE = 27
    MAX_TOLERABLE_SLIDE = 42
    finding_positions = None
    project_name = ""
    project_data = None

    def __init__(self, project, data):
        """ Constructor """
        self.project_name = project
        self.project_data = data
        self.presentation = Presentation(self.template_path)
        self.finding_positions = dict()
        self.finding_positions["solucion_efecto"] = (0,0)
        self.finding_positions["riesgo"] = (1,0)
        self.finding_positions["amenaza"] = (2,0)
        self.finding_positions["vulnerabilidad"] = (3,0)
        self.finding_positions["donde"] = (4,0)
        self.finding_positions["cardinalidad"] = (6,1)
        self.finding_positions["hallazgo"] = (7,0)
        self.finding_positions["categoria"] = (7,1)
        self.finding_positions["criticidad"] = (9,1)
        self.fill_findings_table(data)
        self.fill_results_table(data)
        self.generate_doc()
		
    def get_slide(self, criticality):
        """ 
            Retorna el siguiente slide a editar, dependiendo
            de la criticidad del hallazgo 
        """
        try:
            criticality = float(criticality.replace(",","."))
        except ValueError:
            criticality = 0
        except TypeError:
            criticality = 0
        if criticality == 0:
            return None
        else:
            slide = None
            if criticality > 7: #Critical
                if self.critical_finding_slide < self.MAX_CRITICAL_SLIDE:
                    slide = self.presentation.slides[self.critical_finding_slide]
                    self.critical_finding_slide = self.critical_finding_slide + 1
            elif criticality >= 4 and criticality <= 6.9: #Moderated
                if self.moderate_finding_slide < self.MAX_MODERATE_SLIDE:
                    slide = self.presentation.slides[self.moderate_finding_slide]
                    self.moderate_finding_slide = self.moderate_finding_slide + 1
            else:
                if self.tolerable_finding_slide < self.MAX_TOLERABLE_SLIDE:
                    slide = self.presentation.slides[self.tolerable_finding_slide]
                    self.tolerable_finding_slide = self.tolerable_finding_slide + 1
            return slide

    def fill_findings_table(self, findings):
        """ Llena la tabla de hallazgos resumen del IE Bancolombia"""
        # rows = 2 - 6 vulnerabilidades, cols = v.riesgo c.tecnica nombre 
        finding_table = self.presentation.slides[0].shapes[3].table
        tbl_len = len(findings)
        if len(findings) > 5:
            tbl_len = 5
        init_row = 2
        for i in range(0,tbl_len):
            # 0 = v.riesgo 1 = c.tecnica 2 = nombre 
            finding_table.rows[init_row].cells[0].text_frame.text = ""
            finding_table.rows[init_row].cells[1].text_frame.text = findings[i]["criticidad"]
            finding_table.rows[init_row].cells[2].text_frame.text = findings[i]["hallazgo"]
            init_row += 1

    def fill_results_table(self, findings):
        """ Llena la tabla de resultados resumen del IE Bancolombia """
        # rows = 1-3 hallazgos 4 total , cols = criticidad hallazgos % porcentajes
        results_table = self.presentation.slides[0].shapes[2].table
        crit = 0
        o_crit = 0
        mode = 0
        o_mode = 0
        tole = 0
        o_tole = 0
        total = len(findings)
        for finding in findings:
            criticity = finding["criticidad"]
            ocurrency = finding["cardinalidad"]
            try:
                criticity = float(criticity)
                ocurrency = float(ocurrency)
            except ValueError:
                criticity = 0
                ocurrency = 10000
            if criticity >= 7:
                crit += 1
                o_crit += ocurrency
            elif criticity >= 4 and criticity <= 6.9:
                mode += 1
                o_mode += ocurrency
            else:
                tole += 1
                o_tole += ocurrency
        #hallazgos
        results_table.rows[1].cells[1].text_frame.text = str(crit)
        results_table.rows[2].cells[1].text_frame.text = str(mode)
        results_table.rows[3].cells[1].text_frame.text = str(tole)
        results_table.rows[4].cells[1].text_frame.text = str(total)
        #Porcentaje
        results_table.rows[1].cells[2].text_frame.text = str("%.2f" % float(crit*100/float(total))) + "%"
        results_table.rows[2].cells[2].text_frame.text = str("%.2f" % float(mode*100/float(total))) + "%"
        results_table.rows[3].cells[2].text_frame.text = str("%.2f" % float(tole*100/float(total))) + "%"
        results_table.rows[4].cells[2].text_frame.text = str("%.2f" % 100.0) + "%"
        #Ocurrencias
        results_table.rows[1].cells[3].text_frame.text = str("%d" % o_crit)
        results_table.rows[2].cells[3].text_frame.text = str("%d" % o_mode)
        results_table.rows[3].cells[3].text_frame.text = str("%d" % o_tole)
        results_table.rows[4].cells[3].text_frame.text = str("%d" % int(o_crit + o_mode + o_tole))

    def fill_slide(self, finding):
        """ 
            Asigna a cada figura de texto en el slide
            El valor correspondiente segun el hallazgo 
        """
        try:
            slide = self.get_slide(finding["criticidad"])
            shape, paragraph = self.finding_positions["solucion_efecto"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["solucion_efecto"]
            shape, paragraph = self.finding_positions["riesgo"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["riesgo"]
            shape, paragraph = self.finding_positions["amenaza"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["amenaza"]
            shape, paragraph = self.finding_positions["vulnerabilidad"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["vulnerabilidad"]
            shape, paragraph = self.finding_positions["donde"]
            if len(finding["donde"]) > 200:
                finding["donde"] = finding["donde"][0:200]
            finding["donde"] = finding["donde"].replace("\r", "").replace("\n","; ")
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["donde"]
            shape, paragraph = self.finding_positions["cardinalidad"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["cardinalidad"]
            shape, paragraph = self.finding_positions["hallazgo"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["hallazgo"]
            shape, paragraph = self.finding_positions["criticidad"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["criticidad"]
            #Agregado porque algunos proyectos viejos no tienen categoria en formstack
            if "categoria" in finding:
                shape, paragraph = self.finding_positions["categoria"]
                slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["categoria"]
        except Exception, e:
            return None
		
    def generate_doc(self):
        """ Llena cada slide a partir de un hallazgo """
        for finding in self.project_data:
            self.fill_slide(finding)
        self.save()

    def save(self):
        """ Guarda el documento pptx creado con el nombre del proyecto """
        self.presentation.save(self.result_path.replace(":project",self.project_name))

class Fluid:

    counter = 1
    presentation = None
    template_path = "/var/www/fluid-integrates/app/autodoc/templates/fluid.pptx"
    result_path = "/var/www/fluid-integrates/app/autodoc/results/:project.pptx"
    critical_finding_slide = 1
    moderate_finding_slide = 9
    tolerable_finding_slide = 27
    MAX_CRITICAL_SLIDE = 10
    MAX_MODERATE_SLIDE = 27
    MAX_TOLERABLE_SLIDE = 42
    finding_positions = None
    project_name = ""
    project_data = None

    def __init__(self, project, data):
        """ Constructor"""
        self.project_name = project
        self.project_data = data
        self.presentation = Presentation(self.template_path)
        self.finding_positions = dict()
        self.finding_positions["numero"] = (12,0)
        self.finding_positions["hallazgo"] = (0,0)
        self.finding_positions["criticidad"] = (15,0)
        self.finding_positions["cardinalidad"] = (14,0)
        self.finding_positions["vulnerabilidad"] = (10,0)
        self.finding_positions["vector_ataque"] = (8,0)
        self.finding_positions["amenaza"] = (3,0)
        self.finding_positions["sistema_comprometido"] = (17,0)
        self.finding_positions["solucion_efecto"] = (2,0)
        self.finding_positions["requisitos"] = (6,0)	
        self.fill_findings_table(data)
        self.fill_results_table(data)
        self.generate_doc()
		
    def get_slide(self, criticality):
        """ 
            Retorna el siguiente slide a editar, dependiendo
            de la criticidad del hallazgo 
        """
        try:
            criticality = float(criticality.replace(",","."))
        except ValueError:
            criticality = 0
        except TypeError:
            criticality = 0
        if criticality == 0:
            return None
        else:
            slide = None
            if criticality > 7: #Critical
                if self.critical_finding_slide < self.MAX_CRITICAL_SLIDE:
                    slide = self.presentation.slides[self.critical_finding_slide]
                    self.critical_finding_slide = self.critical_finding_slide + 1
            elif criticality >= 4 and criticality <= 6.9: #Moderated
                if self.moderate_finding_slide < self.MAX_MODERATE_SLIDE:
                    slide = self.presentation.slides[self.moderate_finding_slide]
                    self.moderate_finding_slide = self.moderate_finding_slide + 1
            else:
                if self.tolerable_finding_slide < self.MAX_TOLERABLE_SLIDE:
                    slide = self.presentation.slides[self.tolerable_finding_slide]
                    self.tolerable_finding_slide = self.tolerable_finding_slide + 1
            return slide
    def fill_findings_table(self, findings):
        """ Llena la tabla de hallazgos resumen del IE Fluid"""
        # rows = 2 - 6 vulnerabilidades, cols = criticidad, nombre 
        finding_table = self.presentation.slides[0].shapes[2].table
        tbl_len = len(findings)
        if len(findings) > 5:
            tbl_len = 5
        init_row = 2
        for i in range(0,tbl_len):
            finding_table.rows[init_row].cells[1].text_frame.text = findings[i]["criticidad"]
            finding_table.rows[init_row].cells[2].text_frame.text = findings[i]["hallazgo"]
            init_row += 1

    def fill_results_table(self, findings):
        """ Llena la tabla de resultados resumen del IE Fluid"""
        # rows = 1-3 hallazgos 4 total , cols = criticidad hallazgos % porcentajes
        results_table = self.presentation.slides[0].shapes[1].table
        crit = 0
        o_crit = 0
        mode = 0
        o_mode = 0
        tole = 0
        o_tole = 0
        total = len(findings)
        for finding in findings:
            criticity = finding["criticidad"]
            ocurrency = finding["cardinalidad"]
            try:
                criticity = float(criticity)
                ocurrency = float(ocurrency)
            except ValueError:
                criticity = 0
                ocurrency = 10000
            if criticity >= 7:
                crit += 1
                o_crit += ocurrency
            elif criticity >= 4 and criticity <= 6.9:
                mode += 1
                o_mode += ocurrency
            else:
                tole += 1
                o_tole += ocurrency
        #hallazgos
        results_table.rows[1].cells[1].text_frame.text = str(crit)
        results_table.rows[2].cells[1].text_frame.text = str(mode)
        results_table.rows[3].cells[1].text_frame.text = str(tole)
        results_table.rows[4].cells[1].text_frame.text = str(total)
        #Porcentaje
        results_table.rows[1].cells[2].text_frame.text = str("%.2f" % float(crit*100/float(total))) + "%"
        results_table.rows[2].cells[2].text_frame.text = str("%.2f" % float(mode*100/float(total))) + "%"
        results_table.rows[3].cells[2].text_frame.text = str("%.2f" % float(tole*100/float(total))) + "%"
        results_table.rows[4].cells[2].text_frame.text = str("%.2f" % 100.0) + "%"
        #Ocurrencias
        results_table.rows[1].cells[3].text_frame.text = str("%d" % o_crit)
        results_table.rows[2].cells[3].text_frame.text = str("%d" % o_mode)
        results_table.rows[3].cells[3].text_frame.text = str("%d" % o_tole)
        results_table.rows[4].cells[3].text_frame.text = str("%d" % int(o_crit + o_mode + o_tole))
        
    def fill_slide(self, finding):
        """ 
            Asigna a cada figura de texto en el slide
            El valor correspondiente segun el hallazgo 
        """
        try:
            slide = self.get_slide(finding["criticidad"])
            shape, paragraph = self.finding_positions["numero"]
            current = slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = current.replace("N", str(self.counter))
            shape, paragraph = self.finding_positions["hallazgo"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["hallazgo"]
            shape, paragraph = self.finding_positions["criticidad"]
            current = slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = current.replace("<N>",finding["criticidad"])
            shape, paragraph = self.finding_positions["cardinalidad"]
            current = slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = current.replace("<N>",finding["cardinalidad"])
            shape, paragraph = self.finding_positions["vulnerabilidad"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["vulnerabilidad"]
            shape, paragraph = self.finding_positions["vector_ataque"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["vector_ataque"]
            shape, paragraph = self.finding_positions["amenaza"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["amenaza"]
            shape, paragraph = self.finding_positions["sistema_comprometido"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["sistema_comprometido"]
            shape, paragraph = self.finding_positions["solucion_efecto"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["solucion_efecto"]
            shape, paragraph = self.finding_positions["requisitos"]
            slide.shapes[shape].text_frame.paragraphs[paragraph].runs[0].text = finding["requisitos"]
            self.counter += 1
        except Exception, e:
            return None

    def generate_doc(self):
        """ Llena cada slide a partir de un hallazgo """
        for finding in self.project_data:
            self.fill_slide(finding)
        self.save()

    def save(self):
        """ Guarda el documento pptx creado con el nombre del proyecto """
        self.presentation.save(self.result_path.replace(":project",self.project_name))
